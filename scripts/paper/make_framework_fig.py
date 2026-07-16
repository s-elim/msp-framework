"""Render the MSP method-framework figure (Fig. 1 of the manuscript).

    MPLCONFIGDIR=/tmp python3 scripts/paper/make_framework_fig.py --out paper/figures

Emits fig_framework.png (450 dpi), fig_framework.pdf (vector, for LaTeX), and
fig_framework.pptx (native editable shapes, for talks). All are drawn from the same
layout constants so they cannot drift apart.

The Observation box shows REAL corpus views when paper/figures/framework_views/
holds rgb0.png, rgb1.png, rgb2.png (extracted from a cached LIBERO corpus); it falls
back to schematic frames when they are absent, so the script never fails on a fresh
checkout that has not generated a corpus yet.

Content is the formalization, nothing else: the two learned modules (belief encoder
q_theta, outcome head p_psi), the IB objective L = sum_j beta_j D_j + R with the
sufficiency-for-success budget, and the four inference procedures (certify/abstain,
risk-averse selection, active perception, test-time adaptation) that fall out of the
two modules with no retraining.
"""

from __future__ import annotations

import argparse
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Ellipse, FancyArrowPatch, FancyBboxPatch
from matplotlib.path import Path as MplPath

# ======================================================================================
# Shared layout constants (unit canvas 100 x 64; both renderers scale from here)
# ======================================================================================

CANVAS_W, CANVAS_H = 100.0, 64.0
FIG_W_IN = 7.16  # IEEE two-column text width
FIG_H_IN = FIG_W_IN * CANVAS_H / CANVAS_W

INK = "#22262E"
ARROW = "#3A3F4A"
FAINT = "#6B7280"

C_WORLD_EC, C_WORLD_FC = "#7A8290", "#F4F5F7"
C_OBS_EC, C_OBS_FC = "#5D6B7E", "#EDEFF3"
C_A_EC, C_A_FC, C_A_HD = "#3E6DA6", "#E7EEF7", "#2C4E77"
C_B_EC, C_B_FC, C_B_HD = "#B98A2F", "#FBF1DE", "#7A5A1E"
C_TR_EC, C_TR_FC, C_TR_HD = "#7B5EA7", "#F1ECF7", "#5A4380"
C_INF_EC, C_INF_FC, C_INF_HD = "#3E7D5E", "#EAF2ED", "#2E5F47"
C_TH_EC, C_TH_FC = "#B9BFC9", "#FAFAFB"
C_ABSTAIN = "#A63D3D"
WHITE = "#FFFFFF"

# box: (x, y, w, h) in canvas units, y measured from the bottom
WORLD = (1.0, 40.0, 13.5, 16.0)
OBS = (17.5, 40.0, 12.5, 16.0)
MOD_A = (32.0, 36.5, 20.5, 20.5)
MOD_B = (60.0, 36.5, 15.5, 20.5)
ACTIONS = (58.5, 26.5, 13.0, 6.0)
ORACLE = (1.0, 23.5, 13.5, 10.0)
TRAIN = (17.5, 7.6, 57.5, 11.4)
INFER = (78.0, 7.6, 21.0, 49.4)
THEORY = (1.0, 0.6, 98.0, 3.6)

FLOW_Y = 47.5          # main left-to-right pipeline arrows
GLYPH = (56.25, 51.0)  # belief ellipse centre, in the A-B gap

# feedback arcs as cubic Beziers in two disjoint altitude bands, so they can never
# pinch together: sense is the outer arc (y ~60-62), adapt the inner one (y ~57-59)
SENSE_ARC = ((88.0, 57.0), (88.0, 64.0), (23.75, 64.0), (23.75, 56.2))
ADAPT_ARC = ((80.0, 57.0), (80.0, 60.0), (56.25, 60.0), (56.25, 54.8))

# real-observation thumbnails inside OBS: (x, y) of lower-left corner, side length
THUMBS = [(18.6, 42.2), (20.4, 43.4), (22.2, 44.6)]
THUMB_SIDE = 6.4

# inner rows of the inference column: (key, y, h) within INFER x-extent
INF_ROWS = [
    ("marginalize", 45.4, 7.6),
    ("certify", 36.6, 7.6),
    ("act", 29.4, 6.0),
    ("sense", 22.2, 6.0),
    ("adapt", 15.0, 6.0),
]


# ======================================================================================
# Matplotlib renderer (PNG + PDF)
# ======================================================================================


def _box(ax, rect, fc, ec, lw=0.9, ls="-", radius=0.9, z=2):
    x, y, w, h = rect
    p = FancyBboxPatch(
        (x, y), w, h,
        boxstyle=f"round,pad=0,rounding_size={radius}",
        fc=fc, ec=ec, lw=lw, ls=ls, zorder=z,
    )
    ax.add_patch(p)
    return p


def _arrow(ax, p0, p1, lw=0.9, color=ARROW, ls="-", rad=0.0, z=3, head=4.5):
    a = FancyArrowPatch(
        p0, p1,
        arrowstyle=f"-|>,head_width={head * 0.45},head_length={head}",
        connectionstyle=f"arc3,rad={rad}",
        mutation_scale=1.0, lw=lw, color=color, ls=ls, zorder=z,
        shrinkA=0.0, shrinkB=0.0,
    )
    ax.add_patch(a)


def _elbow(ax, pts, lw=0.9, color=ARROW, ls="-", z=3, head=4.5):
    for i in range(len(pts) - 2):
        ax.plot(*zip(pts[i], pts[i + 1]), lw=lw, color=color, ls=ls, zorder=z,
                solid_capstyle="round")
    _arrow(ax, pts[-2], pts[-1], lw=lw, color=color, ls=ls, z=z, head=head)


def _curve(ax, pts, lw=0.9, color=ARROW, z=3):
    """Cubic-Bezier arrow: pts = (start, ctrl1, ctrl2, end), head at the end."""
    path = MplPath(list(pts), [MplPath.MOVETO, MplPath.CURVE4, MplPath.CURVE4, MplPath.CURVE4])
    a = FancyArrowPatch(path=path, arrowstyle="-|>,head_width=2.0,head_length=4.5",
                        mutation_scale=1.0, lw=lw, color=color, zorder=z)
    ax.add_patch(a)


def _thumb_paths(views_dir: Path) -> list[Path]:
    paths = [views_dir / f"rgb{i}.png" for i in range(3)]
    return paths if all(p.exists() for p in paths) else []


def extract_views(corpus: Path, views_dir: Path) -> None:
    """Pull three real RGB views of one scene out of a cached LIBERO corpus.

    Prefers a curved object (the paper's interesting case). Writes rgb0..rgb2.png at the
    corpus's native resolution; render_* pick them up automatically.
    """
    import torch

    blob = torch.load(corpus, map_location="cpu", weights_only=False)
    obs, names, oi = blob["observation"], blob["object_names"], blob["object_index"]
    want = [i for i, n in enumerate(names)
            if n in ("ketchup", "milk", "orange_juice", "tomato_sauce")]
    scene = int((oi == want[0]).nonzero()[0][0]) if want else 0
    v = obs[scene]  # (V, 4, H, W)
    views_dir.mkdir(parents=True, exist_ok=True)
    for j, vi in enumerate([0, v.shape[0] // 3, (2 * v.shape[0]) // 3]):
        rgb = v[vi, :3].permute(1, 2, 0).float()
        if rgb.max() > 1.5:
            rgb = rgb / 255.0
        plt.imsave(views_dir / f"rgb{j}.png", rgb.clamp(0, 1).numpy())
    print(f"extracted 3 views of scene {scene} ({names[int(oi[scene])]}) -> {views_dir}")


def render_matplotlib(out_dir: Path, views_dir: Path) -> None:
    plt.rcParams.update({
        "font.family": "STIXGeneral",
        "mathtext.fontset": "stix",
        "text.color": INK,
    })
    FS_TITLE, FS_BODY, FS_TINY = 7.6, 6.4, 5.6

    fig = plt.figure(figsize=(FIG_W_IN, FIG_H_IN))
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, CANVAS_W)
    ax.set_ylim(0, CANVAS_H)
    ax.axis("off")

    def T(x, y, s, fs=FS_BODY, ha="center", va="center", w="normal", c=INK, z=4):
        ax.text(x, y, s, fontsize=fs, ha=ha, va=va, fontweight=w, color=c, zorder=z)

    # ---------------------------------------------------------------- world state
    x, y, w, h = WORLD
    _box(ax, WORLD, C_WORLD_FC, C_WORLD_EC, ls=(0, (4, 2.2)))
    T(x + w / 2, y + h - 2.2, "World state", FS_TITLE, w="bold")
    T(x + w / 2, y + h - 4.6, "(hidden)", FS_TINY, c=FAINT)
    T(x + w / 2, y + h - 8.0, r"$x=(S,\,T,\,\phi)$")
    T(x + w / 2, y + h - 11.0, r"shape $\cdot$ pose $\cdot$ physics", FS_TINY, c=FAINT)

    # ---------------------------------------------------------------- observation
    x, y, w, h = OBS
    _box(ax, OBS, C_OBS_FC, C_OBS_EC)
    T(x + w / 2, y + h - 2.2, "Observation", FS_TITLE, w="bold")
    thumbs = _thumb_paths(views_dir)
    if thumbs:
        for j, (tx, ty) in enumerate(THUMBS):
            img = plt.imread(thumbs[j])
            ax.imshow(img, extent=(tx, tx + THUMB_SIDE, ty, ty + THUMB_SIDE),
                      zorder=3 + j, interpolation="bilinear")
            ax.add_patch(FancyBboxPatch((tx, ty), THUMB_SIDE, THUMB_SIDE,
                                        boxstyle="round,pad=0,rounding_size=0.01",
                                        fc="none", ec=WHITE, lw=1.2, zorder=3 + j))
    else:  # fresh checkout without a corpus: schematic frames
        for j, (tx, ty) in enumerate(THUMBS):
            ax.add_patch(FancyBboxPatch((tx, ty), THUMB_SIDE, THUMB_SIDE,
                                        boxstyle="round,pad=0,rounding_size=0.35",
                                        fc=WHITE, ec=C_OBS_EC, lw=0.7, zorder=3 + j))
    T(x + w / 2, y + 1.3, r"$o=\{o_1,\ldots,o_V\}$  RGB-D", FS_TINY)

    # ---------------------------------------------------------------- module A
    x, y, w, h = MOD_A
    _box(ax, MOD_A, C_A_FC, C_A_EC, lw=1.1)
    T(x + w / 2, y + h - 2.2, "Belief encoder", FS_TITLE, w="bold", c=C_A_HD)
    T(x + w / 2, y + h - 4.6, r"$q_\theta(z\,|\,o)$", FS_TITLE, c=C_A_HD)
    stages = [
        r"per-view backbone $f(o_v)$",
        r"permutation-invariant pooling",
        r"$(\mu,\ \log\sigma^2)$",
    ]
    top0 = y + h - 7.2
    for i, s in enumerate(stages):
        by_ = top0 - i * 3.5
        ax.add_patch(FancyBboxPatch((x + 1.8, by_ - 2.6), w - 3.6, 2.6,
                                    boxstyle="round,pad=0,rounding_size=0.5",
                                    fc=WHITE, ec=C_A_EC, lw=0.7, zorder=3))
        T(x + w / 2, by_ - 1.25, s, FS_TINY, z=4)
        if i < 2:
            _arrow(ax, (x + w / 2, by_ - 2.6), (x + w / 2, by_ - 3.45), lw=0.7, head=3.0)
    T(x + w / 2, y + 1.8, r"a belief over $z$, not a pose", FS_TINY, c=C_A_HD)

    # ---------------------------------------------------------------- belief glyph
    bx, by = GLYPH
    for rw, rh, al in ((4.6, 2.6, 0.35), (3.2, 1.8, 0.55), (1.7, 0.95, 0.9)):
        ax.add_patch(Ellipse((bx, by), rw, rh, fc=C_A_FC, ec=C_A_EC, lw=0.7,
                             alpha=al, zorder=4))
    T(bx, by + 2.7, "belief", FS_TINY, c=C_A_HD)
    T(bx, FLOW_Y - 1.7, r"$z_k\sim q_\theta$", FS_TINY)

    # ---------------------------------------------------------------- module B
    x, y, w, h = MOD_B
    _box(ax, MOD_B, C_B_FC, C_B_EC, lw=1.1)
    T(x + w / 2, y + h - 2.2, "Outcome head", FS_TITLE, w="bold", c=C_B_HD)
    T(x + w / 2, y + h - 4.6, r"$p_\psi(y\,|\,z,a,g)$", FS_TITLE, c=C_B_HD)
    ax.add_patch(FancyBboxPatch((x + 1.6, y + h - 9.8), w - 3.2, 2.6,
                                boxstyle="round,pad=0,rounding_size=0.5",
                                fc=WHITE, ec=C_B_EC, lw=0.7, zorder=3))
    T(x + w / 2, y + h - 8.5, r"MLP$(z,\,\mathrm{emb}\,a,\,\mathrm{emb}\,g)$", FS_TINY)
    outs = [
        (r"$succ$ $\cdot$ Bernoulli", y + h - 12.4),
        (r"$margin$ $\cdot$ Gaussian", y + h - 14.7),
        (r"$slip$ $\cdot$ zero-infl. log-normal", y + h - 17.0),
    ]
    for s, oy in outs:
        T(x + w / 2, oy, s, FS_TINY)
    T(x + w / 2, y + 1.8, r"all $K\times N_a$ pairs $(z_k,a_i)$", FS_TINY, c=FAINT)

    # ---------------------------------------------------------------- actions
    x, y, w, h = ACTIONS
    _box(ax, ACTIONS, WHITE, C_B_EC, ls=(0, (4, 2.2)))
    T(x + w / 2, y + h - 2.0, "Candidate grasps", FS_TITLE, w="bold", c=C_B_HD)
    T(x + w / 2, y + h - 4.5, r"$a=(T_g,\,w)\sim\rho$", FS_TINY)

    # ---------------------------------------------------------------- physics oracle
    x, y, w, h = ORACLE
    _box(ax, ORACLE, C_WORLD_FC, C_WORLD_EC, ls=(0, (4, 2.2)))
    T(x + w / 2, y + h - 2.2, "Physics oracle", FS_TITLE, w="bold")
    T(x + w / 2, y + h - 5.0, r"$M(y\,|\,x,a)$")
    T(x + w / 2, y + h - 7.8, "training only", FS_TINY, c=FAINT)

    # ---------------------------------------------------------------- training strip
    x, y, w, h = TRAIN
    _box(ax, TRAIN, C_TR_FC, C_TR_EC, lw=1.1)
    T(x + 1.8, y + h - 2.1, "Training: variational information bottleneck (Eq. 10)",
      FS_TITLE, ha="left", w="bold", c=C_TR_HD)
    T(x + w / 2, y + h - 5.6,
      r"$\mathcal{L}=\sum_j \beta_j D_j + R$,"
      r"   $D_j=\mathbb{E}\left[-\log p_\psi(y_j\,|\,z,a)\right]$,"
      r"   $R=\mathbb{E}_o\,\mathrm{KL}\!\left(q_\theta(z|o)\,\Vert\,\mathcal{N}(0,I)\right)$")
    T(x + w / 2, y + h - 8.9,
      r"$\beta_{succ}=\beta\gg\beta_{margin},\,\beta_{slip}$ (Thm. 3)"
      r"$\;\cdot\;$no reconstruction or pose loss anywhere", FS_TINY, c=C_TR_HD)

    # ---------------------------------------------------------------- inference column
    x, y, w, h = INFER
    _box(ax, INFER, C_INF_FC, C_INF_EC, lw=1.1)
    T(x + w / 2, y + h - 2.1, "Inference  (modules frozen)", FS_TITLE, w="bold", c=C_INF_HD)
    rows = {}
    for name, ry, rh in INF_ROWS:
        ax.add_patch(FancyBboxPatch((x + 1.2, ry), w - 2.4, rh,
                                    boxstyle="round,pad=0,rounding_size=0.6",
                                    fc=WHITE, ec=C_INF_EC, lw=0.8, zorder=3))
        rows[name] = (x + 1.2, ry, w - 2.4, rh)

    def rowT(name, dy, s, fs=FS_TINY, w_="normal", c=INK):
        rx, ry, rw, rh = rows[name]
        T(rx + rw / 2, ry + rh - dy, s, fs, w=w_, c=c)

    rowT("marginalize", 1.6, "marginalize  (Eq. 13-14)", w_="bold", c=C_INF_HD)
    rowT("marginalize", 4.0, r"$s(o,a)=\mathbb{E}_z[\sigma_\psi(z,a)]$")
    rowT("marginalize", 6.3, r"$v(o,a)=\mathrm{Var}_z[\sigma_\psi(z,a)]$")

    rowT("certify", 1.6, "certify: split conformal  (Eq. 22-24)", w_="bold", c=C_INF_HD)
    rowT("certify", 4.0, r"$A_{\mathrm{cert}}=\{a:\ C(o,a)=\{1\}\}$")
    rowT("certify", 6.3, r"$\varnothing\Rightarrow$ ABSTAIN  (coverage $\geq 1-\alpha$)",
         c=C_ABSTAIN)

    rowT("act", 1.6, "act: risk-averse  (Eq. 15)", w_="bold", c=C_INF_HD)
    rowT("act", 4.3, r"$a^{*}=\arg\max_{A_{\mathrm{cert}}}\ s-\lambda\, v$")

    rowT("sense", 1.6, "sense: active perception  (Eq. 16-18)", w_="bold", c=C_INF_HD)
    rowT("sense", 4.3, r"$U(o)>\tau_U\Rightarrow$ acquire view $b^{*}$")

    rowT("adapt", 1.6, r"adapt: after probe $(a_p,y_p)$  (Eq. 19-21)", w_="bold", c=C_INF_HD)
    rowT("adapt", 4.3, r"$q'(z)\propto q_\theta(z|o)\, p_\psi(y_p|z,a_p)$")

    T(x + w / 2, y + 4.4, "no retraining at deployment", FS_TINY, c=C_INF_HD)

    # ---------------------------------------------------------------- theory footer
    x, y, w, h = THEORY
    _box(ax, THEORY, C_TH_FC, C_TH_EC, lw=0.8)
    T(x + 3.4, y + h / 2, "Theory", FS_TITLE, w="bold", c=FAINT)
    T(x + 7.6, y + h / 2,
      r"$Z-O-X-Y$ Markov $\Rightarrow I(Z;Y|A)\leq I(O;Y|A)$ (DPI)"
      r"  $\bullet$  $\eta(o)$ is minimal sufficient and the IB optimum recovers it as"
      r" $\beta\rightarrow\infty$ (Thm. 2-3)"
      r"  $\bullet$  states identifiable only up to $\ker J(x)$ (Thm. 4)",
      FS_TINY, ha="left")

    # ---------------------------------------------------------------- arrows: main flow
    _arrow(ax, (WORLD[0] + WORLD[2], FLOW_Y), (OBS[0], FLOW_Y), lw=1.0)
    T((WORLD[0] + WORLD[2] + OBS[0]) / 2, FLOW_Y + 1.5, r"$p(o|x)$", FS_TINY)
    _arrow(ax, (OBS[0] + OBS[2], FLOW_Y), (MOD_A[0], FLOW_Y), lw=1.0)
    _arrow(ax, (MOD_A[0] + MOD_A[2], FLOW_Y), (MOD_B[0], FLOW_Y), lw=1.0)
    _arrow(ax, (MOD_B[0] + MOD_B[2], FLOW_Y), (INFER[0], FLOW_Y), lw=1.0)
    T((MOD_B[0] + MOD_B[2] + INFER[0]) / 2, FLOW_Y + 1.5, r"$\sigma_\psi$", FS_TINY)

    # world -> oracle (dashed, train time)
    _arrow(ax, (WORLD[0] + WORLD[2] / 2, WORLD[1]),
           (ORACLE[0] + ORACLE[2] / 2, ORACLE[1] + ORACLE[3]), lw=0.8, ls=(0, (4, 2.2)))
    # actions -> module B (straight up)
    _arrow(ax, (65.0, ACTIONS[1] + ACTIONS[3]), (65.0, MOD_B[1]), lw=1.0)
    T(66.3, 34.4, r"$a$", FS_TINY)
    # oracle -> training strip (elbow)
    _elbow(ax, [(4.5, ORACLE[1]), (4.5, 13.3), (TRAIN[0], 13.3)], lw=0.9)
    T(11.0, 14.9, r"$y=(succ,\,margin,\,slip)$", FS_TINY)
    # modules -> training strip (dotted)
    _arrow(ax, (34.8, MOD_A[1]), (34.8, TRAIN[1] + TRAIN[3]), lw=0.7,
           ls=(0, (1.2, 1.6)), head=3.0)
    T(34.0, 21.8, r"rate $R$", FS_TINY, c=C_TR_HD, ha="right")
    _arrow(ax, (73.5, MOD_B[1]), (73.5, TRAIN[1] + TRAIN[3]), lw=0.7,
           ls=(0, (1.2, 1.6)), head=3.0)
    T(72.7, 21.8, r"distortion $D_j$", FS_TINY, c=C_TR_HD, ha="right")

    # ---------------------------------------------------------------- feedback arcs
    _curve(ax, SENSE_ARC, lw=0.9, color=C_INF_HD)
    T(98.8, 61.0, r"sense: new view $o_{V+1}$", FS_TINY, c=C_INF_HD, ha="right")
    _curve(ax, ADAPT_ARC, lw=0.9, color=C_INF_HD)
    T(68.0, 60.15, r"adapt: update belief", FS_TINY, c=C_INF_HD)

    out_dir.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_dir / "fig_framework.png", dpi=450)
    fig.savefig(out_dir / "fig_framework.pdf")
    plt.close(fig)
    print(f"wrote {out_dir}/fig_framework.png and .pdf")


# ======================================================================================
# python-pptx renderer (editable native shapes, 16:9 slide)
# ======================================================================================

# fit the 100 x 64 canvas to the 7.5 in slide height, centred horizontally
PPT_S = 7.5 / CANVAS_H            # inches per canvas unit
PPT_OX = (13.3333 - CANVAS_W * PPT_S) / 2.0
PT_TITLE, PT_BODY, PT_TINY = 12.5, 11.0, 9.5


def render_pptx(out_dir: Path, png_path: Path, views_dir: Path) -> None:
    from lxml import etree
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_LINE_DASH_STYLE
    from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.oxml.ns import qn
    from pptx.util import Inches, Pt

    def rgb(hexstr):
        return RGBColor.from_string(hexstr.lstrip("#"))

    def IX(u):
        return Inches(PPT_OX + u * PPT_S)

    def IY(u_top):
        return Inches((CANVAS_H - u_top) * PPT_S)

    prs = Presentation()
    prs.slide_width = Inches(13.3333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes

    HA = {"center": PP_ALIGN.CENTER, "left": PP_ALIGN.LEFT, "right": PP_ALIGN.RIGHT}

    def _fmt_run(r, text, size, bold, italic, color, script):
        r.text = text
        f = r.font
        f.size = Pt(size)
        f.bold = bold
        f.italic = italic
        f.name = "Times New Roman"
        f.color.rgb = rgb(color)
        if script:
            r._r.get_or_add_rPr().set("baseline", "-25000" if script < 0 else "30000")

    def _para(tf, parts, size, *, bold=False, italic=False, color=INK, ha="center",
              before=0.0, first=False):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.alignment = HA[ha]
        if before:
            p.space_before = Pt(before)
        if isinstance(parts, str):
            parts = [(parts, 0)]
        for text, script in parts:
            _fmt_run(p.add_run(), text, size, bold, italic, color, script)
        return p

    def _rect(rect, fc, ec, *, dash=None, line_w=1.4, radius=0.10, oval=False):
        x, y, w, h = rect
        kind = MSO_SHAPE.OVAL if oval else MSO_SHAPE.ROUNDED_RECTANGLE
        s = shapes.add_shape(kind, IX(x), IY(y + h), Inches(w * PPT_S), Inches(h * PPT_S))
        if not oval:
            try:
                s.adjustments[0] = radius
            except (IndexError, ValueError):
                pass
        s.fill.solid()
        s.fill.fore_color.rgb = rgb(fc)
        s.line.color.rgb = rgb(ec)
        s.line.width = Pt(line_w)
        if dash is not None:
            s.line.dash_style = dash
        s.shadow.inherit = False
        tf = s.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        for m in ("margin_left", "margin_right"):
            setattr(tf, m, Inches(0.04))
        tf.margin_top = Inches(0.05)
        tf.margin_bottom = Inches(0.02)
        return s

    def _label(cx, cy, parts, size, *, color=INK, italic=False, bold=False, ha="center"):
        if isinstance(parts, str):
            parts = [(parts, 0)]
        n = sum(len(t) for t, _ in parts)
        w_in = n * size / 130.0 + 0.25
        h_in = size / 72.0 * 1.7
        left = {"center": IX(cx) - Inches(w_in / 2), "left": IX(cx),
                "right": IX(cx) - Inches(w_in)}[ha]
        w_in = min(w_in, 13.30 - left / 914400)  # keep the box on the slide
        tb = shapes.add_textbox(left, IY(cy) - Inches(h_in / 2), Inches(w_in), Inches(h_in))
        tf = tb.text_frame
        tf.word_wrap = False
        for m in ("margin_left", "margin_right", "margin_top", "margin_bottom"):
            setattr(tf, m, 0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _para(tf, parts, size, bold=bold, italic=italic, color=color, first=True)
        return tb

    def _tail(line_obj):
        ln = line_obj._get_or_add_ln()
        te = etree.SubElement(ln, qn("a:tailEnd"))
        te.set("type", "triangle")
        te.set("w", "med")
        te.set("len", "med")

    def _conn(p0, p1, *, color=ARROW, w=1.4, dash=None, arrow=True):
        c = shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IX(p0[0]), IY(p0[1]),
                                 IX(p1[0]), IY(p1[1]))
        c.line.color.rgb = rgb(color)
        c.line.width = Pt(w)
        if dash is not None:
            c.line.dash_style = dash
        c.shadow.inherit = False
        if arrow:
            _tail(c.line)
        return c

    def _polyline(pts, *, color=ARROW, w=1.4, arrow=True):
        emu = [(int(IX(px)), int(IY(py))) for px, py in pts]
        fb = shapes.build_freeform(emu[0][0], emu[0][1], scale=1.0)
        fb.add_line_segments(emu[1:], close=False)
        s = fb.convert_to_shape()
        s.fill.background()
        s.line.color.rgb = rgb(color)
        s.line.width = Pt(w)
        s.shadow.inherit = False
        if arrow:
            _tail(s.line)
        return s

    def _bezier3(p0, c1, c2, p1, n=32, **kw):
        pts = []
        for i in range(n + 1):
            t = i / n
            u = 1 - t
            x = u**3 * p0[0] + 3 * t * u**2 * c1[0] + 3 * t**2 * u * c2[0] + t**3 * p1[0]
            y = u**3 * p0[1] + 3 * t * u**2 * c1[1] + 3 * t**2 * u * c2[1] + t**3 * p1[1]
            pts.append((x, y))
        return _polyline(pts, **kw)

    DASH = MSO_LINE_DASH_STYLE.DASH
    DOT = MSO_LINE_DASH_STYLE.ROUND_DOT

    # ---------------------------------------------------------------- world state
    s = _rect(WORLD, C_WORLD_FC, C_WORLD_EC, dash=DASH)
    tf = s.text_frame
    _para(tf, "World state", PT_TITLE, bold=True, first=True)
    _para(tf, "(hidden)", PT_TINY, color=FAINT)
    _para(tf, [("x = (S, T, φ)", 0)], PT_BODY, before=6)
    _para(tf, [("shape · pose · physics", 0)], PT_TINY, color=FAINT, before=4)

    # ---------------------------------------------------------------- observation
    x, y, w, h = OBS
    s = _rect(OBS, C_OBS_FC, C_OBS_EC)
    _para(s.text_frame, "Observation", PT_TITLE, bold=True, first=True)
    thumbs = _thumb_paths(views_dir)
    if thumbs:
        for j, (tx, ty) in enumerate(THUMBS):
            pic = shapes.add_picture(str(thumbs[j]), IX(tx), IY(ty + THUMB_SIDE),
                                     Inches(THUMB_SIDE * PPT_S), Inches(THUMB_SIDE * PPT_S))
            pic.line.color.rgb = rgb(WHITE)
            pic.line.width = Pt(1.5)
            pic.shadow.inherit = False
    else:
        for tx, ty in THUMBS:
            _rect((tx, ty, THUMB_SIDE, THUMB_SIDE), WHITE, C_OBS_EC, line_w=1.0, radius=0.18)
    _label(x + w / 2, y + 1.3,
           [("o = {o", 0), ("1", -1), (", …, o", 0), ("V", -1), ("}  RGB-D", 0)], PT_TINY)

    # ---------------------------------------------------------------- module A
    x, y, w, h = MOD_A
    s = _rect(MOD_A, C_A_FC, C_A_EC, line_w=1.8)
    tf = s.text_frame
    _para(tf, "Belief encoder", PT_TITLE, bold=True, color=C_A_HD, first=True)
    _para(tf, [("q", 0), ("θ", -1), ("(z | o)", 0)], PT_TITLE, color=C_A_HD, before=2)
    stages = [
        [("per-view backbone f(o", 0), ("v", -1), (")", 0)],
        [("permutation-invariant pooling", 0)],
        [("(μ, log σ²)", 0)],
    ]
    top0 = y + h - 7.2
    for i, parts in enumerate(stages):
        by_ = top0 - i * 3.5
        sb = _rect((x + 1.8, by_ - 2.6, w - 3.6, 2.6), WHITE, C_A_EC, line_w=1.0, radius=0.25)
        sbtf = sb.text_frame
        sbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
        sbtf.margin_top = 0
        sbtf.margin_bottom = 0
        _para(sbtf, parts, PT_TINY, first=True)
        if i < 2:
            _conn((x + w / 2, by_ - 2.6), (x + w / 2, by_ - 3.4), w=1.0)
    _label(x + w / 2, y + 1.8, "a belief over z, not a pose", PT_TINY, color=C_A_HD)

    # ---------------------------------------------------------------- belief glyph
    bx, by = GLYPH
    for rw_, rh_, fc_ in ((4.6, 2.6, "#EFF4FA"), (3.2, 1.8, "#DCE7F4"), (1.7, 0.95, WHITE)):
        _rect((bx - rw_ / 2, by - rh_ / 2, rw_, rh_), fc_, C_A_EC, line_w=1.0, oval=True)
    _label(bx, by + 2.7, "belief", PT_TINY, color=C_A_HD)
    _label(bx, FLOW_Y - 1.7, [("z", 0), ("k", -1), (" ~ q", 0), ("θ", -1)], PT_TINY)

    # ---------------------------------------------------------------- module B
    x, y, w, h = MOD_B
    s = _rect(MOD_B, C_B_FC, C_B_EC, line_w=1.8)
    tf = s.text_frame
    _para(tf, "Outcome head", PT_TITLE, bold=True, color=C_B_HD, first=True)
    _para(tf, [("p", 0), ("ψ", -1), ("(y | z, a, g)", 0)], PT_TITLE, color=C_B_HD, before=2)
    sb = _rect((x + 1.6, y + h - 9.8, w - 3.2, 2.6), WHITE, C_B_EC, line_w=1.0, radius=0.25)
    sbtf = sb.text_frame
    sbtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    sbtf.margin_top = 0
    _para(sbtf, "MLP(z, emb a, emb g)", PT_TINY, first=True)
    _label(x + w / 2, y + h - 12.4, "succ · Bernoulli", PT_TINY)
    _label(x + w / 2, y + h - 14.7, "margin · Gaussian", PT_TINY)
    _label(x + w / 2, y + h - 17.0, "slip · zero-infl. log-normal", PT_TINY)
    _label(x + w / 2, y + 1.8,
           [("all K × N", 0), ("a", -1), (" pairs (z", 0), ("k", -1), (", a", 0), ("i", -1),
            (")", 0)], PT_TINY, color=FAINT)

    # ---------------------------------------------------------------- actions
    x, y, w, h = ACTIONS
    s = _rect(ACTIONS, WHITE, C_B_EC, dash=DASH)
    tf = s.text_frame
    _para(tf, "Candidate grasps", PT_TITLE, bold=True, color=C_B_HD, first=True)
    _para(tf, [("a = (T", 0), ("g", -1), (", w) ~ ρ", 0)], PT_TINY, before=2)

    # ---------------------------------------------------------------- physics oracle
    x, y, w, h = ORACLE
    s = _rect(ORACLE, C_WORLD_FC, C_WORLD_EC, dash=DASH)
    tf = s.text_frame
    _para(tf, "Physics oracle", PT_TITLE, bold=True, first=True)
    _para(tf, "M(y | x, a)", PT_BODY, before=2)
    _para(tf, "training only", PT_TINY, color=FAINT, before=2)

    # ---------------------------------------------------------------- training strip
    x, y, w, h = TRAIN
    s = _rect(TRAIN, C_TR_FC, C_TR_EC, line_w=1.8)
    tf = s.text_frame
    _para(tf, "Training: variational information bottleneck (Eq. 10)", PT_TITLE,
          bold=True, color=C_TR_HD, ha="left", first=True)
    _para(tf, [("L = Σ", 0), ("j", -1), (" β", 0), ("j", -1), (" D", 0), ("j", -1),
               (" + R,      D", 0), ("j", -1), (" = E[−log p", 0), ("ψ", -1),
               ("(y", 0), ("j", -1), (" | z, a)],      R = E", 0), ("o", -1),
               (" KL(q", 0), ("θ", -1), ("(z|o) ‖ N(0, I))", 0)], PT_BODY, before=6)
    _para(tf, [("β", 0), ("succ", -1), (" = β ≫ β", 0), ("margin", -1), (", β", 0),
               ("slip", -1), (" (Thm. 3)  ·  no reconstruction or pose loss anywhere", 0)],
          PT_TINY, color=C_TR_HD, before=5)

    # ---------------------------------------------------------------- inference column
    x, y, w, h = INFER
    s = _rect(INFER, C_INF_FC, C_INF_EC, line_w=1.8)
    _para(s.text_frame, "Inference  (modules frozen)", PT_TITLE, bold=True, color=C_INF_HD,
          first=True)

    row_text = {
        "marginalize": [
            ("marginalize  (Eq. 13-14)", True, C_INF_HD),
            ([("s(o, a) = E", 0), ("z", -1), ("[σ", 0), ("ψ", -1), ("(z, a)]", 0)], False, INK),
            ([("v(o, a) = Var", 0), ("z", -1), ("[σ", 0), ("ψ", -1), ("(z, a)]", 0)], False, INK),
        ],
        "certify": [
            ("certify: split conformal  (Eq. 22-24)", True, C_INF_HD),
            ([("A", 0), ("cert", -1), (" = {a :  C(o, a) = {1}}", 0)], False, INK),
            ([("∅ ⇒ ABSTAIN  (coverage ≥ 1 − α)", 0)], False, C_ABSTAIN),
        ],
        "act": [
            ("act: risk-averse  (Eq. 15)", True, C_INF_HD),
            ([("a* = argmax", 0), ("A_cert", -1), ("  s − λ v", 0)], False, INK),
        ],
        "sense": [
            ("sense: active perception  (Eq. 16-18)", True, C_INF_HD),
            ([("U(o) > τ", 0), ("U", -1), ("  ⇒  acquire view b*", 0)], False, INK),
        ],
        "adapt": [
            ([("adapt: after probe (a", 0), ("p", -1), (", y", 0), ("p", -1),
              (")  (Eq. 19-21)", 0)], True, C_INF_HD),
            ([("q′(z) ∝ q", 0), ("θ", -1), ("(z|o) p", 0), ("ψ", -1), ("(y", 0), ("p", -1),
              (" | z, a", 0), ("p", -1), (")", 0)], False, INK),
        ],
    }
    for name, ry, rh in INF_ROWS:
        rs = _rect((x + 1.2, ry, w - 2.4, rh), WHITE, C_INF_EC, line_w=1.2, radius=0.15)
        rtf = rs.text_frame
        for i, (parts, bold, color) in enumerate(row_text[name]):
            _para(rtf, parts, PT_TINY, bold=bold, color=color, before=0 if i == 0 else 2,
                  first=(i == 0))
    _label(x + w / 2, y + 4.4, "no retraining at deployment", PT_TINY, color=C_INF_HD)

    # ---------------------------------------------------------------- theory footer
    x, y, w, h = THEORY
    _rect(THEORY, C_TH_FC, C_TH_EC, line_w=1.0)
    _label(x + 3.4, y + h / 2, "Theory", PT_TITLE, bold=True, color=FAINT)
    _label(x + 7.6, y + h / 2,
           [("Z - O - X - Y Markov  ⇒  I(Z; Y|A) ≤ I(O; Y|A)  (DPI)    •    "
             "η(o) is minimal sufficient and the IB optimum recovers it as β → ∞ (Thm. 2-3)"
             "    •    states identifiable only up to ker J(x) (Thm. 4)", 0)],
           PT_TINY, ha="left")

    # ---------------------------------------------------------------- arrows
    _conn((WORLD[0] + WORLD[2], FLOW_Y), (OBS[0], FLOW_Y))
    _label((WORLD[0] + WORLD[2] + OBS[0]) / 2, FLOW_Y + 1.5, "p(o|x)", PT_TINY)
    _conn((OBS[0] + OBS[2], FLOW_Y), (MOD_A[0], FLOW_Y))
    _conn((MOD_A[0] + MOD_A[2], FLOW_Y), (MOD_B[0], FLOW_Y))
    _conn((MOD_B[0] + MOD_B[2], FLOW_Y), (INFER[0], FLOW_Y))
    _label((MOD_B[0] + MOD_B[2] + INFER[0]) / 2, FLOW_Y + 1.5,
           [("σ", 0), ("ψ", -1)], PT_TINY)

    _conn((WORLD[0] + WORLD[2] / 2, WORLD[1]),
          (ORACLE[0] + ORACLE[2] / 2, ORACLE[1] + ORACLE[3]), dash=DASH, w=1.1)
    _conn((65.0, ACTIONS[1] + ACTIONS[3]), (65.0, MOD_B[1]))
    _label(66.3, 34.4, "a", PT_TINY, italic=True)
    _polyline([(4.5, ORACLE[1]), (4.5, 13.3), (TRAIN[0], 13.3)])
    _label(11.0, 14.9, "y = (succ, margin, slip)", PT_TINY)

    _conn((34.8, MOD_A[1]), (34.8, TRAIN[1] + TRAIN[3]), dash=DOT, w=1.0)
    _label(34.0, 21.8, "rate R", PT_TINY, color=C_TR_HD, ha="right")
    _conn((73.5, MOD_B[1]), (73.5, TRAIN[1] + TRAIN[3]), dash=DOT, w=1.0)
    _label(72.7, 21.8, [("distortion D", 0), ("j", -1)], PT_TINY, color=C_TR_HD, ha="right")

    # feedback arcs
    _bezier3(*SENSE_ARC, color=C_INF_HD, w=1.4)
    _label(98.8, 61.0, [("sense: new view o", 0), ("V+1", -1)], PT_TINY,
           color=C_INF_HD, ha="right")
    _bezier3(*ADAPT_ARC, color=C_INF_HD, w=1.4)
    _label(68.0, 60.15, "adapt: update belief", PT_TINY, color=C_INF_HD)

    # ---------------------------------------------------------------- slide 2: render
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    pic_w_in = 7.5 * CANVAS_W / CANVAS_H
    slide2.shapes.add_picture(str(png_path), Inches((13.3333 - pic_w_in) / 2), Inches(0),
                              width=Inches(pic_w_in), height=Inches(7.5))

    out = out_dir / "fig_framework.pptx"
    prs.save(out)
    print(f"wrote {out}")


# ======================================================================================
# CLI
# ======================================================================================

if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", type=Path, default=Path("paper/figures"))
    ap.add_argument("--views", type=Path, default=Path("paper/figures/framework_views"),
                    help="directory holding rgb0..rgb2.png extracted from a LIBERO corpus")
    ap.add_argument("--extract-from-corpus", type=Path, default=None,
                    help="cached LIBERO corpus .pt to (re)extract the Observation views from")
    args = ap.parse_args()
    if args.extract_from_corpus is not None:
        extract_views(args.extract_from_corpus, args.views)
    render_matplotlib(args.out, args.views)
    render_pptx(args.out, args.out / "fig_framework.png", args.views)
